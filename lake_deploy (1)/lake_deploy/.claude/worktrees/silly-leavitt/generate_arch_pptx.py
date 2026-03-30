"""Generate Architecture Current State vs Future State PowerPoint."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
import copy

# ── Palette ─────────────────────────────────────────────────
C_DARK_BLUE  = RGBColor(0x1F, 0x49, 0x7D)
C_MID_BLUE   = RGBColor(0x2E, 0x74, 0xB5)
C_LIGHT_BLUE = RGBColor(0xBD, 0xD7, 0xEE)
C_TEAL       = RGBColor(0x00, 0x70, 0x7F)
C_GREEN      = RGBColor(0x1A, 0x7A, 0x40)
C_LT_GREEN   = RGBColor(0xC6, 0xEF, 0xCE)
C_ORANGE     = RGBColor(0xC5, 0x5A, 0x11)
C_LT_ORANGE  = RGBColor(0xFC, 0xE4, 0xD6)
C_RED        = RGBColor(0xA4, 0x26, 0x2F)
C_LT_RED     = RGBColor(0xF4, 0xCC, 0xCC)
C_YELLOW     = RGBColor(0xFF, 0xC0, 0x00)
C_LT_YELLOW  = RGBColor(0xFF, 0xF2, 0xCC)
C_WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
C_NEAR_WHITE = RGBColor(0xF5, 0xF7, 0xFA)
C_DARK       = RGBColor(0x1A, 0x1A, 0x2E)
C_GREY       = RGBColor(0x60, 0x60, 0x60)
C_LT_GREY    = RGBColor(0xE5, 0xE5, 0xE5)

SW = Inches(13.33)   # slide width  (widescreen 16:9)
SH = Inches(7.5)     # slide height

prs = Presentation()
prs.slide_width  = SW
prs.slide_height = SH

BLANK = prs.slide_layouts[6]   # completely blank

# ── Low-level helpers ────────────────────────────────────────
def add_rect(slide, x, y, w, h, fill=None, line=None, line_w=Pt(0)):
    from pptx.util import Pt
    shape = slide.shapes.add_shape(1, x, y, w, h)  # MSO_SHAPE_TYPE.RECTANGLE=1
    if fill:
        shape.fill.solid(); shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = line_w
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, x, y, w, h,
             size=Pt(14), bold=False, italic=False,
             color=C_DARK, align=PP_ALIGN.LEFT,
             wrap=True, font_name="Calibri"):
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size      = size
    run.font.bold      = bold
    run.font.italic    = italic
    run.font.color.rgb = color
    run.font.name      = font_name
    return txb

def add_bullet_box(slide, items, x, y, w, h,
                   size=Pt(11), color=C_DARK, bullet_color=None,
                   font_name="Calibri", spacing_before=Pt(4)):
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]; first = False
        else:
            p = tf.add_paragraph()
        p.space_before = spacing_before
        run = p.add_run()
        run.text       = ("• " if not item.startswith("  ") else "    ◦ ") + item.lstrip()
        run.font.size  = size
        run.font.color.rgb = bullet_color if bullet_color else color
        run.font.name  = font_name
    return txb

def slide_header(slide, title, subtitle=None,
                 bg=C_DARK_BLUE, fg=C_WHITE):
    """Full-width header bar."""
    add_rect(slide, 0, 0, SW, Inches(1.1), fill=bg)
    add_text(slide, title,
             Inches(0.35), Inches(0.08), Inches(12), Inches(0.6),
             size=Pt(28), bold=True, color=fg, align=PP_ALIGN.LEFT)
    if subtitle:
        add_text(slide, subtitle,
                 Inches(0.35), Inches(0.65), Inches(11), Inches(0.4),
                 size=Pt(13), color=RGBColor(0xBD,0xD7,0xEE),
                 align=PP_ALIGN.LEFT)

def label_tag(slide, text, x, y, w=Inches(1.8), h=Inches(0.32),
              fill=C_MID_BLUE, fg=C_WHITE, size=Pt(10)):
    add_rect(slide, x, y, w, h, fill=fill)
    add_text(slide, text, x, y, w, h,
             size=size, bold=True, color=fg, align=PP_ALIGN.CENTER)

def divider_line(slide, y, color=C_LT_GREY):
    ln = slide.shapes.add_connector(1, Inches(0.3), y, SW-Inches(0.3), y)
    ln.line.color.rgb = color
    ln.line.width = Pt(0.75)

def pill(slide, text, x, y, w, h, fill, fg, size=Pt(10)):
    r = add_rect(slide, x, y, w, h, fill=fill)
    add_text(slide, text, x, y, w, h,
             size=size, bold=True, color=fg, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ═══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)

# Background gradient simulation — two rects
add_rect(sl, 0, 0, SW, SH, fill=C_DARK)
add_rect(sl, 0, 0, SW, Inches(3.8), fill=C_DARK_BLUE)

# Accent stripe
add_rect(sl, 0, Inches(3.7), SW, Inches(0.08), fill=C_YELLOW)

add_text(sl, "Data Lake Architecture",
         Inches(0.8), Inches(0.9), Inches(11.7), Inches(1.1),
         size=Pt(42), bold=True, color=C_WHITE, align=PP_ALIGN.LEFT)
add_text(sl, "Improvement Recommendations",
         Inches(0.8), Inches(1.9), Inches(11.7), Inches(0.8),
         size=Pt(28), color=C_LIGHT_BLUE, align=PP_ALIGN.LEFT)
add_text(sl, "Current State · Gap Analysis · Future State Architecture",
         Inches(0.8), Inches(2.75), Inches(11.7), Inches(0.5),
         size=Pt(15), italic=True, color=C_LT_GREEN, align=PP_ALIGN.LEFT)

# Stats bar
stats = [("109", "Virtual Views"), ("18", "Physical Tables"), ("5", "Data Sources"), ("92%", "Un-Materialised")]
bw = Inches(2.8)
for i,(val,lbl) in enumerate(stats):
    bx = Inches(0.6) + i*bw
    add_rect(sl, bx, Inches(4.1), Inches(2.5), Inches(1.3), fill=RGBColor(0x2A,0x2A,0x4A))
    add_text(sl, val, bx, Inches(4.15), Inches(2.5), Inches(0.7),
             size=Pt(32), bold=True, color=C_YELLOW, align=PP_ALIGN.CENTER)
    add_text(sl, lbl, bx, Inches(4.8), Inches(2.5), Inches(0.4),
             size=Pt(12), color=C_LIGHT_BLUE, align=PP_ALIGN.CENTER)

add_text(sl, "OnPoint Insights LLC  |  Lake Deploy Platform  |  2026-02-20",
         Inches(0.8), Inches(6.9), Inches(11.7), Inches(0.4),
         size=Pt(11), color=C_GREY, align=PP_ALIGN.LEFT)

# ═══════════════════════════════════════════════════════════
# SLIDE 2 — AGENDA
# ═══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, SW, SH, fill=C_NEAR_WHITE)
slide_header(sl, "Agenda", "What We Will Cover Today")

topics = [
    ("01", "Current Architecture Overview",        "Medallion layers, sources, and existing patterns",          C_MID_BLUE),
    ("02", "The Scalability Problem",               "92% virtual views — the core bottleneck at scale",          C_RED),
    ("03", "9 Improvement Areas",                  "Materialization, orchestration, schema, DQ, multi-tenancy",  C_ORANGE),
    ("04", "Future State Architecture",             "Iceberg tables, DAG pipelines, entity spine, unified S3",    C_GREEN),
    ("05", "Implementation Roadmap",                "4 phases: immediate wins → Iceberg → DAG → multi-tenancy",  C_TEAL),
    ("06", "Before & After Comparison",             "Side-by-side current vs future state summary",              C_DARK_BLUE),
]

for i,(num,title,desc,color) in enumerate(topics):
    yy = Inches(1.25) + i * Inches(0.95)
    add_rect(sl, Inches(0.35), yy, Inches(0.55), Inches(0.65), fill=color)
    add_text(sl, num, Inches(0.35), yy+Inches(0.05), Inches(0.55), Inches(0.55),
             size=Pt(18), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_text(sl, title, Inches(1.05), yy+Inches(0.04), Inches(5.5), Inches(0.32),
             size=Pt(14), bold=True, color=C_DARK_BLUE)
    add_text(sl, desc, Inches(1.05), yy+Inches(0.34), Inches(11.5), Inches(0.28),
             size=Pt(11), color=C_GREY)
    divider_line(sl, yy+Inches(0.72), C_LT_GREY)

# ═══════════════════════════════════════════════════════════
# SLIDE 3 — CURRENT ARCHITECTURE OVERVIEW
# ═══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, SW, SH, fill=C_NEAR_WHITE)
slide_header(sl, "Current Architecture Overview", "Medallion Pattern: Raw → Curated → SSOT")

# 5 source boxes
sources = ["Gaiia\n(GraphQL)", "VETRO\n(REST API)", "Platt\n(CSV Export)", "Intacct\n(SOAP/XML)", "Salesforce\n(AppFlow)"]
src_colors = [C_MID_BLUE, C_TEAL, C_GREEN, C_ORANGE, C_DARK_BLUE]
sw2 = Inches(2.0)
for i,(s,c) in enumerate(zip(sources,src_colors)):
    bx = Inches(0.3) + i*(sw2+Inches(0.17))
    add_rect(sl, bx, Inches(1.25), sw2, Inches(0.75), fill=c)
    add_text(sl, s, bx, Inches(1.25), sw2, Inches(0.75),
             size=Pt(11), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

# Arrow down
add_text(sl, "↓ Lambda / Glue / AppFlow / SOAP ingest ↓",
         Inches(0.3), Inches(2.1), Inches(12.7), Inches(0.35),
         size=Pt(11), italic=True, color=C_GREY, align=PP_ALIGN.CENTER)

# Layer boxes
layers = [
    ("RAW LAYER", "S3: JSONL, CSV, ZIP, XML\nGlue Crawler → Athena External Tables\nPartitioned by plan_id/dt (inconsistent naming)",
     C_DARK_BLUE, C_LIGHT_BLUE, Inches(1.5)),
    ("CURATED LAYER ⚠", "109 Virtual Views + 18 Physical Tables\n92% VIRTUAL — No physical storage!\nJSON extraction re-runs every query",
     C_RED, C_LT_RED, Inches(3.0)),
    ("SSOT LAYER", "8 xwalk + 8 dim EXTERNAL Parquet tables\nViews for current-state + exceptions\nEvidence-backed certification (QID + SQL proof)",
     C_GREEN, C_LT_GREEN, Inches(4.5)),
]
for (title,desc,border,bg,yy) in layers:
    add_rect(sl, Inches(0.3), yy, Inches(12.7), Inches(1.2), fill=bg, line=border, line_w=Pt(2))
    add_text(sl, title, Inches(0.45), yy+Inches(0.06), Inches(3.0), Inches(0.4),
             size=Pt(12), bold=True, color=border)
    add_text(sl, desc, Inches(0.45), yy+Inches(0.44), Inches(12.2), Inches(0.72),
             size=Pt(10), color=C_DARK)
    if "⚠" in title:
        add_text(sl, "⚠  SCALABILITY RISK", Inches(9.5), yy+Inches(0.06), Inches(3.3), Inches(0.38),
                 size=Pt(11), bold=True, color=C_RED, align=PP_ALIGN.RIGHT)

add_text(sl, "↓ ↓ ↓ arrow down through layers ↓ ↓ ↓",
         Inches(0.3), Inches(2.42), Inches(12.7), Inches(0.4),
         size=Pt(10), italic=True, color=C_GREY, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════
# SLIDE 4 — THE SCALABILITY PROBLEM
# ═══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, SW, SH, fill=C_DARK)
slide_header(sl, "The Core Problem: 92% Virtual Views", "Every Dashboard Query Scans Raw S3 — Every Time", bg=C_RED)

# Big stat
add_rect(sl, Inches(0.3), Inches(1.25), Inches(3.5), Inches(2.2), fill=RGBColor(0x3A,0x0A,0x0A))
add_text(sl, "92%", Inches(0.3), Inches(1.3), Inches(3.5), Inches(1.3),
         size=Pt(72), bold=True, color=C_YELLOW, align=PP_ALIGN.CENTER)
add_text(sl, "of curated layer\nis VIRTUAL — no storage",
         Inches(0.3), Inches(2.5), Inches(3.5), Inches(0.8),
         size=Pt(12), color=C_WHITE, align=PP_ALIGN.CENTER)

# Query chain
add_rect(sl, Inches(4.0), Inches(1.25), Inches(9.0), Inches(2.2), fill=RGBColor(0x2A,0x2A,0x2A))
add_text(sl, "Every Dashboard Query Executes This Chain:",
         Inches(4.1), Inches(1.28), Inches(8.8), Inches(0.35),
         size=Pt(11), bold=True, color=C_YELLOW)
chain = [
    "Dashboard tile calls  gaiia_subscriptions_services_current  (VIEW)",
    "  → gaiia_billing_subscriptions_current (VIEW)  ← full raw JSONL scan",
    "  → gaiia_products_current              (VIEW)  ← full raw JSONL scan",
    "       → json_parse() + UNNEST + ROW_NUMBER()  ← CPU-intensive",
    "           → 5 concurrent users = 5× parallel full S3 scans",
]
for i,line in enumerate(chain):
    c = C_LT_RED if "scan" in line.lower() or "cpu" in line.lower() or "5×" in line else C_WHITE
    add_text(sl, line, Inches(4.1), Inches(1.65)+i*Inches(0.33),
             Inches(8.8), Inches(0.3), size=Pt(9.5), color=c,
             font_name="Courier New")

# Impact boxes
impacts = [
    ("💰 Cost",    "$5/TB scanned.\nFull JSONL re-scan\nevery query", C_RED),
    ("⏱ Latency", "15–60 seconds\np95 per dashboard\nquery today", C_ORANGE),
    ("📈 Scale",   ">100K rows\ndocumented timeout\n(use LIMIT)", C_YELLOW),
    ("🔁 Reuse",   "Zero result reuse\nacross concurrent\nusers", C_MID_BLUE),
]
for i,(icon,txt,c) in enumerate(impacts):
    bx = Inches(0.3) + i*Inches(3.2)
    add_rect(sl, bx, Inches(3.7), Inches(3.0), Inches(1.7), fill=RGBColor(0x2A,0x2A,0x4A))
    add_text(sl, icon, bx, Inches(3.72), Inches(3.0), Inches(0.45),
             size=Pt(14), bold=True, color=c, align=PP_ALIGN.CENTER)
    add_text(sl, txt, bx, Inches(4.15), Inches(3.0), Inches(1.2),
             size=Pt(10.5), color=C_WHITE, align=PP_ALIGN.CENTER)

add_text(sl, "At 5–10× data growth, this becomes the primary cost and latency bottleneck.",
         Inches(0.3), Inches(5.6), Inches(12.7), Inches(0.4),
         size=Pt(13), bold=True, italic=True, color=C_YELLOW, align=PP_ALIGN.CENTER)

add_text(sl, "Documented in codebase: '⚠ Large queries (>100K rows) may timeout — use LIMIT clauses'",
         Inches(0.3), Inches(6.05), Inches(12.7), Inches(0.35),
         size=Pt(10), color=C_LT_RED, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════
# SLIDE 5 — 9 IMPROVEMENT AREAS (OVERVIEW)
# ═══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, SW, SH, fill=C_NEAR_WHITE)
slide_header(sl, "9 Improvement Areas", "Prioritised by Risk and Business Impact")

improvements = [
    ("1", "Materialise Curated Layer",     "92% virtual → daily Parquet snapshots",              "P0", C_RED),
    ("2", "DAG-Based Orchestration",       "Sequential bash → Step Functions with gates",         "P0", C_RED),
    ("3", "Checkpoint & Recovery",         "File-based → DynamoDB ledger (atomic resume)",        "P0", C_RED),
    ("4", "Schema Drift Detection",        "Silent Glue → Contract files + SNS alerts",           "P1", C_ORANGE),
    ("5", "Incremental Loads",             "Full CTAS rebuild → Iceberg MERGE INTO upserts",      "P1", C_ORANGE),
    ("6", "Multi-Tenancy",                 "Gaiia-only → All sources tenant-partitioned",         "P1", C_ORANGE),
    ("7", "Data Quality Framework",        "Ad hoc exceptions → dq_run_log + active alerts",      "P1", C_ORANGE),
    ("8", "SSOT Crosswalk / Identity",     "Direct source IDs → entity_spine + SCD Type 2",      "P1", C_ORANGE),
    ("9", "IaC & Observability",           "Partial IaC, no cost tracking → 100% IaC + dashboards","P2", C_TEAL),
]

cols = 3; rows_per_col = 3
for i,(num,title,desc,pri,color) in enumerate(improvements):
    col = i % cols
    row = i // cols
    bx = Inches(0.25) + col * Inches(4.35)
    by = Inches(1.3)  + row * Inches(1.9)
    add_rect(sl, bx, by, Inches(4.1), Inches(1.7), fill=C_WHITE, line=color, line_w=Pt(2))
    add_rect(sl, bx, by, Inches(0.45), Inches(1.7), fill=color)
    add_text(sl, num, bx+Inches(0.01), by+Inches(0.55), Inches(0.43), Inches(0.55),
             size=Pt(18), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_rect(sl, bx+Inches(3.45), by+Inches(0.04), Inches(0.6), Inches(0.28), fill=color)
    add_text(sl, pri, bx+Inches(3.45), by+Inches(0.04), Inches(0.6), Inches(0.28),
             size=Pt(9), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_text(sl, title, bx+Inches(0.55), by+Inches(0.1), Inches(3.4), Inches(0.45),
             size=Pt(12), bold=True, color=color)
    add_text(sl, desc, bx+Inches(0.55), by+Inches(0.55), Inches(3.4), Inches(1.0),
             size=Pt(9.5), color=C_DARK)

# ═══════════════════════════════════════════════════════════
# SLIDE 6 — IMPROVEMENT 1: MATERIALISE (DETAIL)
# ═══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, SW, SH, fill=C_NEAR_WHITE)
slide_header(sl, "Improvement 1: Materialise the Curated Layer", "Highest Impact — Immediate Cost & Latency Reduction", bg=C_GREEN)

# LEFT — current
add_rect(sl, Inches(0.2), Inches(1.25), Inches(6.1), Inches(4.5), fill=C_LT_RED, line=C_RED, line_w=Pt(1.5))
add_text(sl, "❌  CURRENT STATE", Inches(0.35), Inches(1.32), Inches(5.8), Inches(0.4),
         size=Pt(13), bold=True, color=C_RED)
add_text(sl, "Every query → raw S3 scan",
         Inches(0.35), Inches(1.75), Inches(5.8), Inches(0.3), size=Pt(11), bold=True, color=C_RED)
chain_cur = [
    "Dashboard query",
    "  → gaiia_subscriptions_current  (VIEW)",
    "      → billing_subscriptions_current  (VIEW)",
    "          → raw JSONL on S3  ← FULL SCAN",
    "      → products_current  (VIEW)",
    "          → raw JSONL on S3  ← FULL SCAN",
    "              json_parse() + UNNEST + ROW_NUMBER()",
    "              Runs every query. 5 users = 5 scans.",
]
for i,line in enumerate(chain_cur):
    c = C_RED if "SCAN" in line or "5 users" in line else RGBColor(0x40,0x10,0x10)
    add_text(sl, line, Inches(0.35), Inches(2.1)+i*Inches(0.28),
             Inches(5.8), Inches(0.26), size=Pt(8.5), color=c, font_name="Courier New")

# metrics current
add_text(sl, "  Scan per query:  ~10–50 GB raw JSONL",   Inches(0.35),Inches(4.43),Inches(5.8),Inches(0.26),size=Pt(9),color=C_RED)
add_text(sl, "  Cost/query:       $0.05–$0.25",           Inches(0.35),Inches(4.69),Inches(5.8),Inches(0.26),size=Pt(9),color=C_RED)
add_text(sl, "  Latency p95:      15–60 seconds",         Inches(0.35),Inches(4.95),Inches(5.8),Inches(0.26),size=Pt(9),color=C_RED)
add_text(sl, "  Concurrent users: 2–3 before contention", Inches(0.35),Inches(5.21),Inches(5.8),Inches(0.26),size=Pt(9),color=C_RED)

# Arrow
add_text(sl, "→", Inches(6.4), Inches(3.1), Inches(0.5), Inches(0.6),
         size=Pt(36), bold=True, color=C_GREEN, align=PP_ALIGN.CENTER)

# RIGHT — future
add_rect(sl, Inches(7.0), Inches(1.25), Inches(6.1), Inches(4.5), fill=C_LT_GREEN, line=C_GREEN, line_w=Pt(1.5))
add_text(sl, "✅  FUTURE STATE", Inches(7.15), Inches(1.32), Inches(5.8), Inches(0.4),
         size=Pt(13), bold=True, color=C_GREEN)
add_text(sl, "Pipeline materialises once → queries read Parquet",
         Inches(7.15), Inches(1.75), Inches(5.8), Inches(0.3), size=Pt(11), bold=True, color=C_GREEN)
chain_fut = [
    "Daily pipeline (runs once at 02:00 UTC):",
    "  JSON extract + UNNEST + ROW_NUMBER()  ← ONCE",
    "  → Write Parquet snapshot (dt=today)   ← Stored!",
    "",
    "Dashboard query (all day long):",
    "  gaiia_subscriptions_current  (VIEW)",
    "  → SELECT * FROM snapshot WHERE dt=MAX(dt)",
    "  → Reads ~10 MB Parquet  ← FAST & CHEAP",
    "  5 users share same Parquet cache",
]
for i,line in enumerate(chain_fut):
    c = C_GREEN if "ONCE" in line or "FAST" in line or "Stored" in line else RGBColor(0x0A,0x30,0x15)
    add_text(sl, line, Inches(7.15), Inches(2.1)+i*Inches(0.28),
             Inches(5.8), Inches(0.26), size=Pt(8.5), color=c, font_name="Courier New")

add_text(sl, "  Scan per query:  ~10–100 MB Parquet",    Inches(7.15),Inches(4.43),Inches(5.8),Inches(0.26),size=Pt(9),color=C_GREEN)
add_text(sl, "  Cost/query:       <$0.001",               Inches(7.15),Inches(4.69),Inches(5.8),Inches(0.26),size=Pt(9),color=C_GREEN)
add_text(sl, "  Latency p95:      1–3 seconds",           Inches(7.15),Inches(4.95),Inches(5.8),Inches(0.26),size=Pt(9),color=C_GREEN)
add_text(sl, "  Concurrent users: 50+ (shared Parquet)",  Inches(7.15),Inches(5.21),Inches(5.8),Inches(0.26),size=Pt(9),color=C_GREEN)

add_text(sl, "Also: Enable Athena Query Result Reuse (60-min TTL) — zero code change, immediate benefit",
         Inches(0.2), Inches(5.85), Inches(12.9), Inches(0.35),
         size=Pt(10), bold=True, color=C_DARK_BLUE, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════
# SLIDE 7 — IMPROVEMENT 2+3: ORCHESTRATION & CHECKPOINTS
# ═══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, SW, SH, fill=C_NEAR_WHITE)
slide_header(sl, "Improvements 2 & 3: Orchestration + Checkpointing", "From Sequential Bash to DAG-Driven with Atomic Recovery", bg=C_ORANGE)

# Current
add_rect(sl, Inches(0.2), Inches(1.25), Inches(5.9), Inches(3.5), fill=C_LT_RED, line=C_RED, line_w=Pt(1.5))
add_text(sl, "❌  CURRENT — Sequential + Fragile", Inches(0.35), Inches(1.32), Inches(5.6), Inches(0.38),
         size=Pt(12), bold=True, color=C_RED)
cur_orch = [
    "EventBridge (time-based) → Lambda",
    "  → bash: run Glue crawler",
    "  → bash: run CTAS (even if crawler failed!)",
    "  → bash: write manifest",
    "",
    "Checkpointing: plan_index.json on S3",
    "  If file lost → restart from beginning",
    "  Intacct: full restart on any failure",
    "  Rate-limit (429): poll next_allowed_ts",
]
for i,line in enumerate(cur_orch):
    c = C_RED if "even if" in line or "restart" in line or "poll" in line else RGBColor(0x40,0x10,0x10)
    add_text(sl, line, Inches(0.35), Inches(1.75)+i*Inches(0.29),
             Inches(5.6), Inches(0.27), size=Pt(9), color=c, font_name="Courier New")

# Future DAG
add_rect(sl, Inches(6.3), Inches(1.25), Inches(6.8), Inches(3.5), fill=C_LT_GREEN, line=C_GREEN, line_w=Pt(1.5))
add_text(sl, "✅  FUTURE — DAG + DynamoDB Checkpoints", Inches(6.45), Inches(1.32), Inches(6.5), Inches(0.38),
         size=Pt(12), bold=True, color=C_GREEN)

dag_steps = [
    ("START 02:00 UTC", C_DARK_BLUE),
    ("[parallel] Gaiia / Vetro / Platt / Intacct / SF ingest", C_MID_BLUE),
    ("↓ GATE: all partitions exist + freshness OK", C_ORANGE),
    ("[parallel] Materialise curated Parquet snapshots", C_MID_BLUE),
    ("↓ GATE: DQ exception_rate < threshold", C_ORANGE),
    ("[parallel] SSOT certification queries", C_GREEN),
    ("Write manifest + update ops dashboard", C_GREEN),
    ("END — target < 45 minutes total", C_DARK_BLUE),
]
for i,(step,c) in enumerate(dag_steps):
    add_text(sl, step, Inches(6.45), Inches(1.75)+i*Inches(0.36),
             Inches(6.5), Inches(0.33), size=Pt(9.5),
             bold=("GATE" in step or "START" in step or "END" in step),
             color=c, font_name="Courier New")

# Checkpoint comparison
add_rect(sl, Inches(0.2), Inches(4.95), Inches(5.9), Inches(1.7), fill=C_LT_ORANGE, line=C_ORANGE, line_w=Pt(1))
add_text(sl, "❌  plan_index.json  — S3 file-based checkpoint",
         Inches(0.35), Inches(5.0), Inches(5.6), Inches(0.33), size=Pt(10), bold=True, color=C_ORANGE)
add_text(sl, "File lost = full restart from beginning\n429 handling = poll every invocation",
         Inches(0.35), Inches(5.35), Inches(5.6), Inches(0.9), size=Pt(10), color=C_DARK)

add_rect(sl, Inches(6.3), Inches(4.95), Inches(6.8), Inches(1.7), fill=C_LT_GREEN, line=C_GREEN, line_w=Pt(1))
add_text(sl, "✅  DynamoDB Checkpoint Ledger  — atomic",
         Inches(6.45), Inches(5.0), Inches(6.5), Inches(0.33), size=Pt(10), bold=True, color=C_GREEN)
add_text(sl, "Atomic updates — resume from last cursor on restart\n429 → SQS delay queue (event-driven, no polling)",
         Inches(6.45), Inches(5.35), Inches(6.5), Inches(0.9), size=Pt(10), color=C_DARK)

add_text(sl, "→", Inches(6.0), Inches(3.1), Inches(0.4), Inches(0.5),
         size=Pt(30), bold=True, color=C_GREEN, align=PP_ALIGN.CENTER)
add_text(sl, "→", Inches(6.0), Inches(5.35), Inches(0.4), Inches(0.5),
         size=Pt(30), bold=True, color=C_GREEN, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════
# SLIDE 8 — IMPROVEMENTS 4–7 SUMMARY
# ═══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, SW, SH, fill=C_NEAR_WHITE)
slide_header(sl, "Improvements 4–7: Schema · DQ · Multi-Tenancy · Iceberg", "Four Cross-Cutting Platform Enhancements")

panels = [
    ("4", "Schema Drift Detection", C_TEAL,
     ["Current: Glue detects new fields — no alert",
      "Future:  Schema contract files per entity",
      "         Glue → Lambda → SNS on any change",
      "         Null-rate probe after every ingest",
      "         schema_version on all views"]),
    ("5", "Incremental Loads → Iceberg", C_MID_BLUE,
     ["Current: Full daily CTAS rebuild — entire table",
      "Future:  Apache Iceberg tables in Athena",
      "         MERGE INTO for upserts (no rebuild)",
      "         Time-travel queries (audit history)",
      "         Auto-compaction of small files"]),
    ("6", "Full Multi-Tenancy", C_GREEN,
     ["Current: Only Gaiia has tenant= partitions",
      "Future:  ALL sources: tenant/source/entity/dt=",
      "         tenants.json config → zero code change",
      "         New operator = config entry, not code",
      "         Cross-tenant analytics in one query"]),
    ("7", "Data Quality Framework", C_ORANGE,
     ["Current: Exceptions captured — not monitored",
      "Future:  curated_recon.dq_run_log (central)",
      "         config/dq_thresholds.json (versioned)",
      "         CloudWatch alarm: exception_rate > X",
      "         Quarantine prefix for bad-row isolation"]),
]

for i,(num,title,color,bullets) in enumerate(panels):
    col = i % 2; row = i // 2
    bx = Inches(0.2) + col * Inches(6.55)
    by = Inches(1.25) + row * Inches(2.85)
    add_rect(sl, bx, by, Inches(6.35), Inches(2.6), fill=C_WHITE, line=color, line_w=Pt(2))
    add_rect(sl, bx, by, Inches(6.35), Inches(0.45), fill=color)
    add_text(sl, f"  {num}  {title}", bx+Inches(0.05), by+Inches(0.04),
             Inches(6.2), Inches(0.38), size=Pt(13), bold=True, color=C_WHITE)
    for j,b in enumerate(bullets):
        c = (C_RED if b.startswith("Current") else
             C_GREEN if b.startswith("Future") else C_DARK)
        add_text(sl, b, bx+Inches(0.15), by+Inches(0.55)+j*Inches(0.38),
                 Inches(6.05), Inches(0.36), size=Pt(9.5), color=c)

# ═══════════════════════════════════════════════════════════
# SLIDE 9 — IMPROVEMENT 8+9: SSOT SPINE & IaC
# ═══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, SW, SH, fill=C_NEAR_WHITE)
slide_header(sl, "Improvements 8 & 9: Identity Spine + IaC & Observability", "Structural and Operational Completeness")

# Left — Identity spine
add_rect(sl, Inches(0.2), Inches(1.25), Inches(6.2), Inches(4.9), fill=C_WHITE, line=C_ORANGE, line_w=Pt(2))
add_rect(sl, Inches(0.2), Inches(1.25), Inches(6.2), Inches(0.45), fill=C_ORANGE)
add_text(sl, "  8  SSOT Entity Spine + SCD Type 2", Inches(0.35), Inches(1.29),
         Inches(5.9), Inches(0.38), size=Pt(13), bold=True, color=C_WHITE)

add_text(sl, "❌  Current: Source IDs used directly — no surrogate key\n"
              "        SF Account ID, Platt ID, Intacct ID used interchangeably\n"
              "        SF → Intacct crosswalk is non-1:1 (migrations create M:M)\n"
              "        Customer changes silently overwrite — no history",
         Inches(0.35), Inches(1.8), Inches(5.9), Inches(1.2), size=Pt(9.5), color=C_RED)

add_text(sl, "✅  Future: entity_spine Iceberg table\n"
              "        ssot_entity_id (UUID) — single key for all sources\n"
              "        sf_id | platt_id | intacct_id | gaiia_id columns\n"
              "        confidence: high | medium | low | unresolved\n"
              "        SCD Type 2: valid_from, valid_to, is_current",
         Inches(0.35), Inches(3.1), Inches(5.9), Inches(1.4), size=Pt(9.5), color=C_GREEN)

add_text(sl, "✅  All views JOIN through entity_spine — never directly to source IDs\n"
              "✅  M:M crosswalk resolved at identity layer — not at query time\n"
              "✅  Historical customer changes auditable via time-travel",
         Inches(0.35), Inches(4.6), Inches(5.9), Inches(0.9), size=Pt(9.5), color=C_DARK_BLUE)

# Right — IaC & Observability
add_rect(sl, Inches(6.6), Inches(1.25), Inches(6.5), Inches(4.9), fill=C_WHITE, line=C_TEAL, line_w=Pt(2))
add_rect(sl, Inches(6.6), Inches(1.25), Inches(6.5), Inches(0.45), fill=C_TEAL)
add_text(sl, "  9  IaC Completeness + Observability", Inches(6.75), Inches(1.29),
         Inches(6.2), Inches(0.38), size=Pt(13), bold=True, color=C_WHITE)

iac_rows = [
    ("Glue Crawlers",      "Console only",    "CloudFormation"),
    ("AppFlow",            "Console only",    "CF ConnectorProfile"),
    ("S3 Lifecycle Rules", "Console only",    "IaC — raw→Glacier 90d"),
    ("Athena Workgroups",  "Console only",    "IaC + cost by stage"),
    ("CI/CD Auth",         "Long-lived keys", "OIDC short-lived STS"),
    ("Frontend Deploy",    "Manual upload",   "GitHub → Amplify auto"),
]
add_text(sl, "Component", Inches(6.75), Inches(1.8), Inches(1.8), Inches(0.3),
         size=Pt(9), bold=True, color=C_TEAL)
add_text(sl, "Current", Inches(8.6), Inches(1.8), Inches(1.8), Inches(0.3),
         size=Pt(9), bold=True, color=C_RED)
add_text(sl, "Future", Inches(10.45), Inches(1.8), Inches(2.5), Inches(0.3),
         size=Pt(9), bold=True, color=C_GREEN)
for i,(comp,cur,fut) in enumerate(iac_rows):
    y = Inches(2.15) + i * Inches(0.38)
    add_text(sl, comp, Inches(6.75), y, Inches(1.8), Inches(0.35), size=Pt(9), color=C_DARK)
    add_text(sl, "❌ " + cur, Inches(8.6),  y, Inches(1.8), Inches(0.35), size=Pt(9), color=C_RED)
    add_text(sl, "✅ " + fut, Inches(10.45), y, Inches(2.5), Inches(0.35), size=Pt(9), color=C_GREEN)

add_text(sl, "Observability Additions:",
         Inches(6.75), Inches(4.6), Inches(6.2), Inches(0.3), size=Pt(10), bold=True, color=C_TEAL)
obs = [
    "CloudWatch Alarm: max_dt age > 26h per source (freshness SLA)",
    "Alarm: exception_rate > threshold → SNS / PagerDuty",
    "Athena cost by workgroup: raw | curated | ssot stages",
    "Pipeline P95 duration metric + 2× baseline alarm",
]
for i,o in enumerate(obs):
    add_text(sl, "• " + o, Inches(6.75), Inches(4.95)+i*Inches(0.3),
             Inches(6.2), Inches(0.28), size=Pt(9), color=C_DARK)

# ═══════════════════════════════════════════════════════════
# SLIDE 10 — FUTURE STATE ARCHITECTURE
# ═══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, SW, SH, fill=C_DARK)
slide_header(sl, "Future State: Full Architecture", "Materialised Medallion with DAG, Iceberg & Tenant Partitioning", bg=C_GREEN)

# Source tier
add_text(sl, "DATA SOURCES", Inches(0.25), Inches(1.25), Inches(12.8), Inches(0.3),
         size=Pt(9), bold=True, color=C_LT_GREY, align=PP_ALIGN.CENTER)
sources2 = [("Gaiia\nGraphQL",C_MID_BLUE),("VETRO\nREST",C_TEAL),("Platt\nCSV→Parquet",C_GREEN),
            ("Intacct\nSOAP",C_ORANGE),("Salesforce\nAppFlow",C_DARK_BLUE)]
for i,(s,c) in enumerate(sources2):
    bx = Inches(0.25) + i*Inches(2.6)
    add_rect(sl, bx, Inches(1.6), Inches(2.4), Inches(0.65), fill=c)
    add_text(sl, s, bx, Inches(1.6), Inches(2.4), Inches(0.65),
             size=Pt(10), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

add_text(sl, "⬇  Step Functions DAG — parallel ingest with dependency gates  ⬇",
         Inches(0.25), Inches(2.35), Inches(12.8), Inches(0.3),
         size=Pt(10), italic=True, color=C_YELLOW, align=PP_ALIGN.CENTER)

# Raw tier
add_rect(sl, Inches(0.2), Inches(2.72), Inches(12.9), Inches(0.85), fill=RGBColor(0x1A,0x2A,0x4A), line=C_MID_BLUE, line_w=Pt(1))
add_text(sl, "RAW LAYER  —  tenant/source/entity/dt=YYYY-MM-DD/  |  NDJSON.GZ standard  |  Partition projection enabled  |  Immutable append-only",
         Inches(0.35), Inches(2.78), Inches(12.6), Inches(0.65),
         size=Pt(10), color=C_LIGHT_BLUE, align=PP_ALIGN.CENTER)

add_text(sl, "⬇  Materialisation step: JSON extract + UNNEST + ROW_NUMBER() runs ONCE daily  ⬇",
         Inches(0.25), Inches(3.65), Inches(12.8), Inches(0.28),
         size=Pt(9.5), italic=True, color=C_LT_GREEN, align=PP_ALIGN.CENTER)

# Curated tier
add_rect(sl, Inches(0.2), Inches(4.0), Inches(12.9), Inches(0.95), fill=RGBColor(0x0A,0x2A,0x1A), line=C_GREEN, line_w=Pt(1.5))
add_text(sl, "CURATED LAYER  —  Parquet snapshots (ZSTD, 128–256 MB)  |  tenant/entity/dt=  |  Thin views on top  |  Athena Query Result Reuse (60-min TTL)",
         Inches(0.35), Inches(4.07), Inches(12.6), Inches(0.75),
         size=Pt(10), color=C_LT_GREEN, align=PP_ALIGN.CENTER)

add_text(sl, "⬇  DQ gate: exception_rate checked via dq_run_log  ⬇",
         Inches(0.25), Inches(5.02), Inches(12.8), Inches(0.25),
         size=Pt(9.5), italic=True, color=C_YELLOW, align=PP_ALIGN.CENTER)

# SSOT tier
add_rect(sl, Inches(0.2), Inches(5.33), Inches(12.9), Inches(0.85), fill=RGBColor(0x0A,0x1A,0x3A), line=C_YELLOW, line_w=Pt(1.5))
add_text(sl, "SSOT LAYER  —  Apache Iceberg tables (MERGE INTO, time-travel, auto-compaction)  |  entity_spine (SCD Type 2)  |  SSOT certification with QID proof",
         Inches(0.35), Inches(5.38), Inches(12.6), Inches(0.75),
         size=Pt(10), color=C_YELLOW, align=PP_ALIGN.CENTER)

# DQ + Observability sidebar note
add_text(sl, "DynamoDB Checkpoints  |  Schema Contract Files  |  CloudWatch Alarms (freshness + DQ + DLQ)  |  Athena cost by workgroup",
         Inches(0.25), Inches(6.3), Inches(12.8), Inches(0.3),
         size=Pt(9), color=C_GREY, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════
# SLIDE 11 — IMPLEMENTATION ROADMAP
# ═══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, SW, SH, fill=C_NEAR_WHITE)
slide_header(sl, "Implementation Roadmap", "4 Phases — Quick Wins to Full Transformation")

phases = [
    ("Phase 0", "Week 1\nQuick Wins", C_GREEN, [
        "Enable Athena Query Result Reuse — 30 min",
        "CloudWatch Alarm on Vetro SQS DLQ",
        "Externalise DQ thresholds to JSON config",
        "Add WHERE dt=MAX(dt) guard to *_current views",
        "Standardise partition key (dt everywhere)",
        "Add tenant column to all curated views",
    ]),
    ("Phase 1", "Weeks 2–4\nMaterialization", C_ORANGE, [
        "Materialise top 5 curated entities as daily Parquet",
        "DynamoDB checkpoint ledger (replace plan_index.json)",
        "Create dq_run_log table + populate daily",
        "CloudWatch Alarm: max_dt age > 26h per source",
        "Schema contract files for Gaiia + Vetro entities",
        "Codify Glue crawlers in CloudFormation",
    ]),
    ("Phase 2", "Month 2\nIceberg & DAG", C_TEAL, [
        "Migrate high-churn entities to Apache Iceberg",
        "Build Step Functions DAG (replaces sequential bash)",
        "entity_spine surrogate key table (Iceberg)",
        "SCD Type 2 for customer/account history",
        "Athena Workgroups by stage + cost dashboards",
        "OIDC GitHub→AWS auth in CI pipeline",
    ]),
    ("Phase 3", "Month 3+\nMulti-Tenant", C_DARK_BLUE, [
        "Unified S3 layout: tenant/source/entity/dt= (all sources)",
        "Parameterise Lambda/Glue with tenant_id from config",
        "Data freshness + DQ dashboards in CloudWatch",
        "dbt tests or Great Expectations on curated tables",
        "AppFlow → IaC (CF ConnectorProfile)",
        "Amplify + GitHub CI/CD auto-deploy on merge",
    ]),
]

for i,(phase,period,color,items) in enumerate(phases):
    bx = Inches(0.2) + i * Inches(3.27)
    by = Inches(1.25)
    add_rect(sl, bx, by, Inches(3.1), Inches(5.8), fill=C_WHITE, line=color, line_w=Pt(2))
    add_rect(sl, bx, by, Inches(3.1), Inches(1.0), fill=color)
    add_text(sl, phase, bx+Inches(0.1), by+Inches(0.05), Inches(2.9), Inches(0.45),
             size=Pt(16), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_text(sl, period, bx+Inches(0.1), by+Inches(0.5), Inches(2.9), Inches(0.48),
             size=Pt(10), color=C_WHITE, align=PP_ALIGN.CENTER)
    for j,item in enumerate(items):
        add_text(sl, "• " + item, bx+Inches(0.12), by+Inches(1.1)+j*Inches(0.73),
                 Inches(2.88), Inches(0.65), size=Pt(9), color=C_DARK)

# Timeline bar
add_rect(sl, Inches(0.2), Inches(7.1), Inches(12.9), Inches(0.2), fill=C_LT_GREY)
for i,color in enumerate([C_GREEN, C_ORANGE, C_TEAL, C_DARK_BLUE]):
    add_rect(sl, Inches(0.2)+i*Inches(3.27), Inches(7.1), Inches(3.1), Inches(0.2), fill=color)

# ═══════════════════════════════════════════════════════════
# SLIDE 12 — BEFORE & AFTER SCORECARD
# ═══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, SW, SH, fill=C_NEAR_WHITE)
slide_header(sl, "Before & After: Architecture Scorecard", "Side-by-Side Comparison of All 9 Improvement Areas")

rows_data = [
    ("Curated Layer",     "92% virtual views — raw scan per query",  "Daily Parquet snapshots + Iceberg for SSOT"),
    ("Orchestration",     "Sequential bash + EventBridge time-based", "Step Functions DAG with dependency gates"),
    ("Checkpointing",     "S3 file (plan_index.json) — full restart", "DynamoDB ledger — atomic, cursor-based resume"),
    ("Schema Mgmt",       "Glue detects drift — no alert",           "Contract files + Lambda → SNS on drift"),
    ("Incremental Loads", "Full daily CTAS rebuild",                  "Iceberg MERGE INTO — upserts only"),
    ("Multi-Tenancy",     "Gaiia only (tenant= partition)",           "All sources tenant/source/entity/dt="),
    ("Data Quality",      "Exceptions captured — not alarmed",        "dq_run_log + CloudWatch alarms + quarantine"),
    ("SSOT Identity",     "Source IDs directly — non-1:1 xwalk",      "entity_spine UUID + SCD Type 2"),
    ("IaC & Observability","Partial IaC, no cost tracking, no alarms","100% IaC + freshness/cost/DQ dashboards"),
]

hdr_y = Inches(1.25)
add_rect(sl, Inches(0.2), hdr_y, Inches(2.5), Inches(0.4), fill=C_DARK_BLUE)
add_rect(sl, Inches(2.75), hdr_y, Inches(4.8), Inches(0.4), fill=C_RED)
add_rect(sl, Inches(7.6), hdr_y, Inches(5.5), Inches(0.4), fill=C_GREEN)
add_text(sl, "Domain", Inches(0.2), hdr_y, Inches(2.5), Inches(0.4),
         size=Pt(10), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
add_text(sl, "❌  Current State", Inches(2.75), hdr_y, Inches(4.8), Inches(0.4),
         size=Pt(10), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
add_text(sl, "✅  Future State", Inches(7.6), hdr_y, Inches(5.5), Inches(0.4),
         size=Pt(10), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

for i,(domain,cur,fut) in enumerate(rows_data):
    y = Inches(1.7) + i * Inches(0.56)
    fill = 'FFFFFF' if i % 2 == 0 else 'F5F7FA'
    for bx,bw,color in [(Inches(0.2), Inches(2.5), None),
                         (Inches(2.75), Inches(4.8), None),
                         (Inches(7.6), Inches(5.5), None)]:
        r = sl.shapes.add_shape(1, bx, y, bw, Inches(0.52))
        r.fill.solid(); r.fill.fore_color.rgb = RGBColor(*[int(fill[i:i+2],16) for i in (0,2,4)])
        r.line.color.rgb = C_LT_GREY; r.line.width = Pt(0.5)

    add_text(sl, domain, Inches(0.25), y+Inches(0.08), Inches(2.4), Inches(0.4),
             size=Pt(9.5), bold=True, color=C_DARK_BLUE)
    add_text(sl, cur, Inches(2.8), y+Inches(0.06), Inches(4.7), Inches(0.42),
             size=Pt(9), color=C_RED)
    add_text(sl, fut, Inches(7.65), y+Inches(0.06), Inches(5.4), Inches(0.42),
             size=Pt(9), color=C_GREEN)

# ═══════════════════════════════════════════════════════════
# SLIDE 13 — KEY METRICS: IMPACT SUMMARY
# ═══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, SW, SH, fill=C_DARK)
slide_header(sl, "Expected Impact: Key Metrics", "Before vs After Materialisation + DAG + Iceberg", bg=C_MID_BLUE)

metrics = [
    ("Athena Scan\nper Query",    "~10–50 GB",     "~10–100 MB",   "500× reduction",  C_GREEN),
    ("Query Cost",                "$0.05–$0.25",   "<$0.001",      "50–250× cheaper", C_GREEN),
    ("Dashboard\nLatency (p95)",  "15–60 sec",     "1–3 sec",      "20× faster",      C_GREEN),
    ("Concurrent\nUsers",         "2–3 (contention)","50+ (shared Parquet)", "20× more", C_GREEN),
    ("JSON Parsing\nOverhead",    "Every query",   "Once per day", "Full elimination", C_YELLOW),
    ("Pipeline\nRecovery",        "Full restart",  "Resume from cursor", "Zero data loss", C_YELLOW),
    ("Schema\nBreaking Changes",  "Silent failure", "SNS alert",   "Immediate detect", C_YELLOW),
    ("IaC Coverage",              "~60%",          "100%",         "Full auditability",C_TEAL),
]

for i,(label,before,after,gain,color) in enumerate(metrics):
    col = i % 4; row = i // 4
    bx = Inches(0.2) + col * Inches(3.27)
    by = Inches(1.25) + row * Inches(2.7)
    add_rect(sl, bx, by, Inches(3.1), Inches(2.5), fill=RGBColor(0x1E,0x1E,0x3E))
    add_text(sl, label, bx+Inches(0.1), by+Inches(0.1), Inches(2.9), Inches(0.55),
             size=Pt(11), bold=True, color=color, align=PP_ALIGN.CENTER)
    add_text(sl, "Before", bx+Inches(0.1), by+Inches(0.72), Inches(2.9), Inches(0.25),
             size=Pt(8), color=C_GREY, align=PP_ALIGN.CENTER)
    add_text(sl, before, bx+Inches(0.1), by+Inches(0.95), Inches(2.9), Inches(0.38),
             size=Pt(11), color=C_LT_RED, align=PP_ALIGN.CENTER)
    add_text(sl, "After", bx+Inches(0.1), by+Inches(1.38), Inches(2.9), Inches(0.25),
             size=Pt(8), color=C_GREY, align=PP_ALIGN.CENTER)
    add_text(sl, after, bx+Inches(0.1), by+Inches(1.6), Inches(2.9), Inches(0.38),
             size=Pt(11), color=C_LT_GREEN, align=PP_ALIGN.CENTER)
    add_rect(sl, bx+Inches(0.3), by+Inches(2.05), Inches(2.5), Inches(0.3), fill=color)
    add_text(sl, gain, bx+Inches(0.3), by+Inches(2.05), Inches(2.5), Inches(0.3),
             size=Pt(9), bold=True, color=C_DARK, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════
# SLIDE 14 — CLOSING / NEXT STEPS
# ═══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, SW, SH, fill=C_DARK_BLUE)
add_rect(sl, 0, 0, SW, Inches(0.08), fill=C_YELLOW)
add_rect(sl, 0, SH-Inches(0.08), SW, Inches(0.08), fill=C_YELLOW)

add_text(sl, "Next Steps", Inches(0.8), Inches(0.5), Inches(11.7), Inches(0.8),
         size=Pt(38), bold=True, color=C_WHITE, align=PP_ALIGN.LEFT)

steps = [
    ("This Week", C_GREEN, [
        "Enable Athena Query Result Reuse on dashboard workgroup",
        "Add CloudWatch Alarm on Vetro SQS DLQ",
        "Externalise DQ thresholds to config/dq_thresholds.json",
    ]),
    ("Weeks 2–4", C_YELLOW, [
        "Materialise top 5 curated entities as Parquet snapshots",
        "Deploy DynamoDB checkpoint ledger",
        "Create dq_run_log and wire to CloudWatch alarms",
    ]),
    ("Month 2", C_LIGHT_BLUE, [
        "Migrate to Apache Iceberg for SSOT tables",
        "Build Step Functions DAG replacing sequential bash",
        "Deploy entity_spine with SCD Type 2",
    ]),
]

for i,(period,color,items) in enumerate(steps):
    bx = Inches(0.4) + i * Inches(4.27)
    by = Inches(1.55)
    add_rect(sl, bx, by, Inches(4.0), Inches(0.45), fill=color)
    add_text(sl, period, bx, by, Inches(4.0), Inches(0.45),
             size=Pt(14), bold=True, color=C_DARK, align=PP_ALIGN.CENTER)
    for j,item in enumerate(items):
        add_text(sl, "▸  " + item, bx+Inches(0.1), by+Inches(0.55)+j*Inches(0.65),
                 Inches(3.8), Inches(0.58), size=Pt(11), color=C_WHITE)

add_text(sl, "The architecture is well-designed. These improvements secure its scalability for 5–10× growth.",
         Inches(0.6), Inches(4.7), Inches(12.1), Inches(0.5),
         size=Pt(14), italic=True, color=C_LT_GREEN, align=PP_ALIGN.CENTER)

add_text(sl, "OnPoint Insights LLC  ·  Lake Deploy Platform  ·  2026-02-20  ·  Confidential",
         Inches(0.6), Inches(7.05), Inches(12.1), Inches(0.3),
         size=Pt(9), color=RGBColor(0x80,0x90,0xA0), align=PP_ALIGN.CENTER)

# ── SAVE ────────────────────────────────────────────────────
out = ("/Users/vinaymistry/Library/CloudStorage/OneDrive-OnPointInsightsLLC/"
       "GitRepo/lake_deploy/.claude/worktrees/silly-leavitt/Architecture_Improvement_Presentation.pptx")
prs.save(out)
print(f"PowerPoint saved: {out}")
